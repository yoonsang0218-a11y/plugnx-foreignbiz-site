# -*- coding: utf-8 -*-
"""외국인 사업 Appendix 2장 — 경영전략회의 덱 톤 맞춤(블루 eyebrow·네이비 박스·▶ 블루)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLUE  = RGBColor(0x1A, 0x4F, 0xD8)
NAVY  = RGBColor(0x1E, 0x2B, 0x57)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x70, 0x76, 0x83)
LIGHT = RGBColor(0xEE, 0xF1, 0xF6)
LINE  = RGBColor(0xD7, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]
ML = 0.55  # margin left
CW = 13.333 - ML * 2


def _ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def run(p, text, size, color, bold=False, spacing=None):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = FONT
    f.color.rgb = color
    _ea(r, FONT)
    if spacing is not None:
        r._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
    return r


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def rect(slide, l, t, w, h, fill, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def hline(slide, l, t, w, color=LINE, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(l), Inches(t), Inches(l + w), Inches(t))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln


def header(slide, eyebrow, headline):
    tf = box(slide, ML, 0.42, CW, 0.35)
    run(tf.paragraphs[0], eyebrow, 11.5, BLUE, bold=True, spacing=0.6)
    tf2 = box(slide, ML, 0.82, CW, 1.05)
    p = tf2.paragraphs[0]; p.line_spacing = 1.05
    run(p, headline, 23, INK, bold=True)
    hline(slide, ML, 1.95, CW, LINE, 1.0)


def navybox(slide, t, label, summary, label_w=1.15, h=0.92):
    rect(slide, ML, t, label_w, h, NAVY, rounded=True)
    tfl = box(slide, ML, t, label_w, h)
    tfl.vertical_anchor = MSO_ANCHOR.MIDDLE
    pl = tfl.paragraphs[0]; pl.alignment = PP_ALIGN.CENTER
    run(pl, label, 15, WHITE, bold=True)
    rect(slide, ML + label_w + 0.18, t, CW - label_w - 0.18, h, LIGHT, rounded=True)
    tfs = box(slide, ML + label_w + 0.42, t + 0.05, CW - label_w - 0.62, h - 0.1)
    tfs.vertical_anchor = MSO_ANCHOR.MIDDLE
    ps = tfs.paragraphs[0]; ps.line_spacing = 1.18
    run(ps, summary, 12.5, INK)


def takeaway(slide, t, text):
    hline(slide, ML, t, CW, LINE, 1.0)
    tf = box(slide, ML, t + 0.12, CW, 0.6)
    p = tf.paragraphs[0]; p.line_spacing = 1.1
    run(p, "▶  ", 13, BLUE, bold=True)
    run(p, text, 13, BLUE, bold=True)


def footer(slide, page):
    tf = box(slide, ML, 7.08, CW - 1.2, 0.3)
    run(tf.paragraphs[0],
        "참고용·법적효력 없음. 재무수치는 목표·가정이며 PoC·LOI 실측 대상. 출처: 법무부 출입국·외국인정책본부·고용노동부·금융위·사업기획서 v1.0(SSOT).",
        8, MUTED)
    tfp = box(slide, 13.333 - 1.5, 7.05, 1.0, 0.3)
    pp = tfp.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    run(pp, page, 9.5, MUTED)
    tfb = box(slide, 13.333 - 1.5, 6.78, 1.0, 0.28)
    pb = tfb.paragraphs[0]; pb.alignment = PP_ALIGN.RIGHT
    rr = run(pb, "PLUG", 10, NAVY, bold=True); run(pb, "nx", 10, BLUE, bold=True)


def bullet(tf, text, size=11.5, color=INK, first=False, gap=6, lead="· "):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.line_spacing = 1.15; p.space_after = Pt(gap)
    run(p, lead, size, BLUE if lead.strip() else color, bold=True)
    run(p, text, size, color)
    return p


# ============ SLIDE 1 — 추진배경 ============
s = prs.slides.add_slide(BL)
header(s, "Appendix: 외국인 재직관리·환급 사업 — ① 추진배경",
       "정부 정책이 '들어온 외국인을 오래·합법으로 붙드는 인프라'로 이동 — '채용 이후'는 무주공산")
navybox(s, 2.12, "배경",
        "체류외국인 278.3만 시대 — 기업은 엑셀로 관리하고, 근로자는 받을 돈을 놓치며, 금융사는 외국인을 평가 못 한다. "
        "세 통증을 같은 데이터로 푸는 빈 공간. 당사 인증(ForeignID)·금융권(케이뱅크 등) 관계가 그대로 진입 자산.")

blocks = [
    ("① 세 통증 동시 발생",
     ["기업 — 비자만료·자격외취업 시 사업주 제재",
      "근로자 — 미수령 환급·금융 접근 부재",
      "금융사 — 외국인 thin-file(신용정보 부재)"]),
    ("② 정책 골든윈도",
     ["기조 「양적 확대」 → 「수급 안정+관리·신뢰」",
      "E-7-4 숙련전환 점수제(107개 지역) 〔법무부〕",
      "기본인적정보 ICAO 표준·비금융 대안신용 확대 〔금융위〕"]),
    ("③ 빈 공간 (무주공산)",
     ["'채용 이후 통합관리+환급+양면데이터' 보유사 없음",
      "경쟁사는 입국 전·유학생·채용입구·생활정보에 분산",
      "→ 당사가 메울 공백"]),
]
bw = (CW - 0.6) / 3
by = 3.35
for i, (title, items) in enumerate(blocks):
    bx = ML + i * (bw + 0.3)
    rect(s, bx, by, bw, 2.75, WHITE, line=LINE, rounded=True)
    tf = box(s, bx + 0.25, by + 0.22, bw - 0.5, 0.45)
    run(tf.paragraphs[0], title, 13.5, NAVY, bold=True)
    tfb = box(s, bx + 0.25, by + 0.78, bw - 0.5, 1.8)
    for j, it in enumerate(items):
        bullet(tfb, it, size=11, first=(j == 0))
takeaway(s, 6.32, "기업 관리부담↑ + 근로자 권익공백 + 금융 thin-file이 동시에 커진다 → 같은 데이터 하나로 푸는 인프라 기회")
footer(s, "1 / 2")

# ============ SLIDE 2 — 시장규모 ============
s = prs.slides.add_slide(BL)
header(s, "Appendix: 외국인 재직관리·환급 사업 — ② 시장규모",
       "Q1. 시장 규모는 괜찮은가?  (체류외국인 278.3만 · 취업자격 59.4만)")
navybox(s, 2.12, "A1.",
        "B2B TAM 49만(E계열 취업비자) · B2C TAM 216만(장기체류). 외국인 은행계좌 696만인데 인터넷은행 신용대출 ≈0 — 데이터 공백이 곧 기회.")


def table(slide, l, t, w, title, rows):
    tf = box(slide, l, t, w, 0.32)
    run(tf.paragraphs[0], title, 12.5, NAVY, bold=True)
    gt = slide.shapes.add_table(len(rows) + 1, 3, Inches(l), Inches(t + 0.4),
                                Inches(w), Inches(0.42 * (len(rows) + 1))).table
    gt.columns[0].width = Inches(w * 0.30)
    gt.columns[1].width = Inches(w * 0.50)
    gt.columns[2].width = Inches(w * 0.20)
    hdr = ["구분", "값", "라벨"]
    for c in range(3):
        cell = gt.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        run(p, hdr[c], 9.5, WHITE, bold=True)
    for r_i, rowv in enumerate(rows, start=1):
        for c in range(3):
            cell = gt.cell(r_i, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r_i % 2 else LIGHT
            cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.line_spacing = 1.0
            run(p, rowv[c], 9.5, INK if c < 2 else MUTED, bold=(c == 0))


tw = (CW - 0.4) / 2
table(s, ML, 3.35, tw, "시장 모수", [
    ["체류외국인 총계", "2,783,247명", "법무부 '25.12"],
    ["취업자격", "594,047명", "법무부 '25.12"],
    ["B2B TAM (VisaDesk)", "49만 (E계열)", "SSOT"],
    ["B2C TAM (BackWon)", "216만 (장기체류)", "SSOT"],
])
table(s, ML + tw + 0.4, 3.35, tw, "구성 · 금융 공백", [
    ["비자별", "E-9 33.7만·H-2 9.3만·E-7 6.4만·F-4 55.6만", "'24말 연보"],
    ["국적", "중국 35.2%·베트남 12.1%", "법무부"],
    ["금융 공백", "계좌 696만 vs 인터넷은행 여신 ≈0", "'26.4"],
    ["목표(가정)", "PoC 15사 3,300만→ARR 21억→60억", "SSOT"],
])
takeaway(s, 6.32, "E-9 13만→8만으로 신규유입은 줄지만, 정책 무게가 '존량(재직·전환·정착) 관리'로 이동 → 우리 타깃은 오히려 두터워짐")
footer(s, "2 / 2")

out = r"output/plugnx-외국인-추진배경-시장규모-matched.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
